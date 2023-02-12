import os
from datetime import date
from time import sleep
from typing import Callable

import dash_bootstrap_components as dbc
import dash_uploader as du
import diskcache
import pandas as pd
from dash import Dash, Input, Output, State, ctx, dash, dcc
from dash.long_callback import DiskcacheLongCallbackManager
from plotly.graph_objects import Figure
from pptx import Presentation

import database.driver as driver
import vis.prs_lib as prsLib
from computation.data import DataPings, DataSessions, LicenseUsage
from computation.features import Features
from vis.additional_data_vis import (
    get_cas_statistics,
    get_license_usage_table,
    get_multi_total_amount_table,
    get_package_combination_table,
    get_total_amount_table,
)
from vis.graph_vis import (
    empty_fig,
    get_cas_graph,
    get_cluster_id_comparison_graph,
    get_fpc_graph,
    get_multi_cas_graph,
    get_multi_files_graph,
    get_token_graph,
)
from vis.web_designs import DROPDOWN_OPTIONS, tab_layout

GRAPH_LINE_COLOR = "#FFFFFF"
UPLOAD_CACHE_PATH = os.path.abspath("./cache/upload_data/")
HIGH_PERF_MODE = True

# Diskcache - needed for long_callbacks
cache = diskcache.Cache(os.path.abspath("./cache"))
long_callback_manager = DiskcacheLongCallbackManager(cache)

# Dash
app = Dash(
    external_stylesheets=[dbc.themes.BOOTSTRAP],
    long_callback_manager=long_callback_manager,
)
app.title = "Token Dashboard Threedy"
app._favicon = "threedy_favicon.png"
app.layout = dbc.Container(tab_layout(), fluid=True)

du.configure_upload(app, UPLOAD_CACHE_PATH, use_upload_id=False)


def convert_report_to_df(name: str):
    """
    Parameters
    ----------
    name : String
        the name of a csv-file

    Returns
    -------
    pd.Dataframe which represents the csv-file
    """
    return pd.read_csv(UPLOAD_CACHE_PATH + "/" + name)


def select_graph(
    menu_entry: str, session: DataSessions, identifier: str, graph_type: str
):
    """
    Parameters
    ----------
    menu_entry : String which defines the selected graph
    session : pd.Dataframe which represents the datapoints used for the selected graph
    identifier : str which represents the identifier of the currently chosen file
    graph_type: either "bar" or "line"

    Returns
    -------
    plotly.express figure which is a selected graphical representation of the dataframe
    pd.DataFrame : DataFrame ready for dbc.Table component
    """
    fig = empty_fig()
    additional = pd.DataFrame()
    idents = driver.get_df_from_db("identifier")
    idents = (
        idents[idents["Type"] == "Feature"]["FileIdentifier"]
        .to_numpy(dtype=str)
        .flatten()
    )

    if menu_entry == "Token Consumption":
        fig = get_token_graph(session, graph_type=graph_type)
        additional = get_total_amount_table(session)
        fig.update_layout(
            xaxis_title="Time",
            yaxis_title="Token",
            legend_title="Products",
        )

    elif menu_entry == "Product Usage":
        fig = get_fpc_graph(session)
        fig.update_layout(xaxis_title="Products", yaxis_title="Usage (%)")
        additional = get_package_combination_table(session)

    elif menu_entry == "Concurrent Active Sessions":
        fig = get_cas_graph(session, graph_type=graph_type)
        fig.update_layout(xaxis_title="Time", yaxis_title="CAS")
        additional = get_cas_statistics(session)
    elif menu_entry == "Cluster-ID Comparison":
        c_ids = driver.get_df_from_db("cluster_ids")
        c_ids = c_ids[c_ids["identifier"] == identifier]
        c_ids = c_ids["cluster_id"].to_numpy().tolist()
        fig = get_cluster_id_comparison_graph(session, c_ids)
        fig.update_layout(
            xaxis_title="Time", yaxis_title="Token", legend_title="Idents"
        )
    elif menu_entry == "File Comparison (Token)":
        fig = get_multi_files_graph(session, idents, graph_type=graph_type)
        fig.update_layout(
            xaxis_title="Time", yaxis_title="Token", legend_title="Idents"
        )
        additional = get_multi_total_amount_table(session, idents)
    elif menu_entry == "File Comparison (CAS)":
        fig = get_multi_cas_graph(session, idents)
        fig.update_layout(xaxis_title="Time", yaxis_title="CAS", legend_title="Idents")

    if HIGH_PERF_MODE:
        fig.update_traces(hovertemplate=None, hoverinfo="skip")

    fig.update_xaxes(
        showline=True,
        linecolor=GRAPH_LINE_COLOR,
        gridcolor=GRAPH_LINE_COLOR,
        zerolinecolor=GRAPH_LINE_COLOR,
        zerolinewidth=1,
    )
    fig.update_yaxes(
        showline=True,
        linecolor=GRAPH_LINE_COLOR,
        gridcolor=GRAPH_LINE_COLOR,
        zerolinecolor=GRAPH_LINE_COLOR,
        zerolinewidth=1,
    )
    fig.update_layout(
        font_family="sans-serif",
        font_color=GRAPH_LINE_COLOR,
        plot_bgcolor="rgba(0,0,0,0)",
        margin=dict(l=70, r=15, t=15, b=15),
    )

    return fig, additional


def select_date(sel_date, df: pd.DataFrame, asc: bool, init_change: bool):
    """
    Parameters
    ----------
    sel_date : String which represents a date, selected in the calendar tool
    df : pd.Dataframe used to extract a date out of this dataframe
    asc : boolean which is true if the date should be the beginning date and
            false if the date should be the ending date
    init_change : boolean which is true if a new date should be loaded out of the dataframe

    Returns
    -------
    date.Date which represents a given date
    """
    if sel_date is None or init_change:
        data = df.sort_values(
            by="block_start", ascending=asc, ignore_index=True
        ).block_start[0]
    else:
        data = sel_date
        data = data.split("T")
        data = data[0].split()
        data = data[0].split("-")
        data = date(int(data[0]), int(data[1]), int(data[2]))
    return data


@app.callback(
    Output(component_id="graph_data3", component_property="children"),
    Output("license-store", "data"),
    Input("file-select-license", "value"),
    prevent_inital_call=True,
)
def update_output_license(filename: str):
    """
    Parameters
    ----------
    filename: String

    main computation of the frontend

    Returns
    -------
    dcc.Graph which represents the graph with the id 'graph3'
    """
    if driver.check_if_table_exists("license"):
        license_data = driver.get_df_from_db("license")
        license_data = license_data[license_data["identifier"] == filename]
        license_usage = LicenseUsage(license_data)
        additional = get_license_usage_table(license_usage)
        return dbc.Table.from_dataframe(additional), additional.to_dict()
    else:
        return dash.no_update, dash.no_update


@app.callback(
    Output(component_id="overview_filename", component_property="children"),
    Output(component_id="overview_lines", component_property="children"),
    Output(component_id="overview_cal_days", component_property="children"),
    Output(component_id="overview_metered_days", component_property="children"),
    Output(component_id="graphs-store", component_property="children"),
    Output("additionals-store", "data"),
    Output("select-date", "start_date"),
    Output("select-date", "end_date"),
    Input("select-date", "start_date"),
    Input("select-date", "end_date"),
    Input("filename", "data"),
    Input(component_id="graph-type", component_property="value"),
    Input("file-select", "value"),
    Input("cluster_id-select", "value"),
    prevent_inital_call=True,
)
def update_output_div(
    start_date: str,
    end_date: str,
    filename: str,
    graph_type: str,
    file_select: str,
    c_id_select: str,
):
    """
    Parameters
    ----------
    start_date : String which represents the selected entry of the beginning day in the calendar tool
                    with the id 'select-date'
    end_date : String which represents the selected entry of the ending day in the calendar tool
                    with the id 'select-date'
    filename : str
    graph_type : String, either "bar", "line" or "auto"
    file_select: String
    c_id_select: String

    main computation of the frontend

    Returns
    -------
    String which the number of days inbetween the selected days of the calendar tool
            or the first and last day of the dataframe
    String which represents the metered days in the dataframe
    dcc.Graph which represents the graph with the id 'graph1'
    dcc.Graph which represents the graph with the id 'graph2'
    String which represents additional data to the graph with the id 'graph1'
    String which represents additional data to the graph with the id 'graph2'
    date.Date which represents the first selected day in the calendar tool
    date.Date which represents the last selected day in the calendar tool
    """

    if driver.check_if_table_exists("session"):
        features = Features().get_data_features()

        """computation if the data is loaded first or the selected date is changed"""
        if (
            (not driver.check_if_table_exists("current_data"))
            or ctx.triggered_id == "select-date"
            or ctx.triggered_id == "file-select"
            or ctx.triggered_id == "cluster_id-select"
            or ctx.triggered_id == "graph-type"
        ):
            """converting the dict containing all data back into a pd.Dataframe"""
            sql_session = driver.get_df_from_db("session")

            """converting Strings representing date into data.Date dates"""
            if c_id_select == "All Cluster-IDs":
                c_id = None
            else:
                c_id = c_id_select
            data_pings = DataPings(
                filename, driver.get_df_from_db("pings"), features, c_id
            )
            sessions = DataSessions(
                sql_session, data_pings, features, 300, file_select, c_id
            )

            """checking if new data is loaded and new initial dates should be set"""
            new_data = False
            if not driver.check_if_table_exists("current_data"):
                new_data = True

            """setting the first and last dates for the represented data"""
            last = select_date(end_date, sessions.data, False, new_data)
            first = select_date(start_date, sessions.data, True, new_data)

            """trimming the dataframe to fit into the selected dates"""
            sessions.crop_data(first, last)

            if graph_type == "automatic":
                graph_type = (
                    "bar" if (len(data_pings.get_metered_days()) <= 2) else "line"
                )

            graphs = []
            additionals = dict()
            for option in DROPDOWN_OPTIONS:
                id = option["label"]
                fig, additional = select_graph(id, sessions, file_select, graph_type)
                graphs.append(dcc.Graph(figure=fig, className="graph"))
                additionals[str(option["value"])] = additional.to_dict()

            driver.df_to_sql_replace(sessions.data, "current_data")

            filename, lines, cal_days, metered_days = get_overview_table(
                sessions.data_pings, file_select, c_id
            )

            return (
                filename,
                lines,
                cal_days,
                metered_days,  # overview_table data
                graphs,
                additionals,  # graphs-store, additionals-store
                first,
                last,
            )

        """converting the dict containing all data back into a pd.Dataframe"""
        df_current = driver.get_df_from_db("current_data")

        if c_id_select == "All Cluster-IDs":
            c_id = None
        else:
            c_id = c_id_select
        data_pings = DataPings(filename, driver.get_df_from_db("pings"), features, c_id)
        sessions = DataSessions(
            df_current, data_pings, features, 300, file_select, c_id
        )

    filename, lines, cal_days, metered_days = get_overview_table(None, "", None)

    return (
        filename,
        lines,
        cal_days,
        metered_days,  # overview_table data
        dash.no_update,
        dash.no_update,  # graphs-store, additionals-store
        start_date,
        end_date,
    )


@app.callback(
    Output(component_id="graph1", component_property="children"),
    Output(component_id="graph2", component_property="children"),
    Output(component_id="graph_data1", component_property="children"),
    Output(component_id="graph_data2", component_property="children"),
    Input("graphs-store", "children"),
    Input("additionals-store", "data"),
    Input(component_id="dropdown1", component_property="value"),
    Input(component_id="dropdown2", component_property="value"),
    prevent_inital_call=True,
)
def update_dropdown(
    graphs: list,
    additionals: dict,
    drop1: int,
    drop2: int,
):
    graph1 = graphs[drop1]
    additional1 = dbc.Table.from_dataframe(
        pd.DataFrame.from_dict(additionals[str(drop1)])
    )

    graph2 = graphs[drop2]
    additional2 = dbc.Table.from_dataframe(
        pd.DataFrame.from_dict(additionals[str(drop2)])
    )

    return graph1, graph2, additional1, additional2


def get_overview_table(data_pings: DataPings, file_identifier: str, cluster_id: str):
    """
    Get overview table data.

    Parameters
    ----------
    data_pings : DataPings which represents the selected entry in the dropdown menu with the id 'dropdown1'
    file_identifier: String which represents the identifier of each file
    cluster_id: String which represents the cluster_id that is selected

    Returns
    -------
    filename: str
        name of report
    lines: str
        number of pings
    cal_days: str
        number of calendar days
    metered_days: str
        number of metered days
    """
    if data_pings is not None:
        data_pings_new = DataPings(
            file_identifier,
            data_pings.data[data_pings.data["identifier"] == file_identifier],
            data_pings.features,
            cluster_id,
        )
        filename = str(data_pings_new.get_filename())
        lines = str(len(data_pings_new.get_pings().index))
        cal_days = str(len(data_pings_new.get_sequence_of_days()))
        metered_days = str(len(data_pings_new.get_metered_days()))
    else:
        filename = ""
        lines = ""
        cal_days = ""
        metered_days = ""

    return filename, lines, cal_days, metered_days


@app.long_callback(
    output=[
        Output("confirm", "n_clicks"),
        Output("dash-uploader", "isCompleted"),
        Output("filename", "data"),
        Output("filename_license", "data"),
    ],
    inputs=[
        Input("dash-uploader", "isCompleted"),
        State("dash-uploader", "fileNames"),
        Input("confirm", "n_clicks"),
    ],
    running=[
        (Output("main", "style"), {"filter": "grayscale(100%)"}, {}),
        (Output("header", "style"), {"filter": "grayscale(100%)"}, {}),
        (
            Output("progress_div", "style"),
            {"visibility": "visible"},
            {"visibility": "hidden"},
        ),
    ],
    progress=[
        Output("progress_bar", "value"),
        Output("progress_bar", "label"),
        Output("progress_bar_header", "children"),
        Output("progress_message", "children"),
        Output("modal_ident", "is_open"),
    ],
    prevent_inital_call=True,
)
def load_data(set_progress: Callable, is_com: bool, files: str, confirm: int):
    """
    loads the data into the database

    Parameters
    ----------
    set_progress : Function to update the website while this callback function is running
    is_com : bool
        which indicates if the upload is completed
    files : String
        array of filenames (only first element used)
    confirm : int
        number of time the confirm button is clicked

    Returns
    -------
    int which represents the number of lines in the csv-file
    String which represents the filename of the loaded csv-file
    bool which indicates if the pop-up window is open (True) or closed (False)
    Resets the clicks counter of the confirm button
    Resets the upload
    """
    if is_com:
        header_text = "Upload Report"
        set_progress((100, "0/5", header_text, "Waiting for input", True))

        while confirm is None:
            sleep(0.1)

        # 1. Convert Data
        set_progress((0, "0/5", header_text, "Converting Data", False))
        # without "sleep()" the user cannot see the progress on the website
        sleep(1)

        filename = files[0].split(".")[0]

        datagram = convert_report_to_df(files[0])

        if "grant_id" in datagram.columns:
            set_progress((60, "3/5", header_text, "Loading License Data", False))
            sleep(1)
            ident = driver.get_last_input("identifier")
            table = driver.get_df_from_db("identifier")
            table.loc[len(table.index) - 1, "Type"] = "License"
            driver.df_to_sql_replace(table, "identifier")

            license_data = datagram
            license_data["identifier"] = ident
            driver.df_to_sql_append(datagram, "license")
            set_progress((100, "5/5", header_text, "Loaded Data Successfully", False))
            sleep(1)
            return None, False, dash.no_update, filename

        set_progress(
            (20, "1/5", header_text, "Getting Number of Lines and Features", False)
        )
        sleep(1)

        # 2. Get Features
        features = Features().get_data_features()
        set_progress((40, "2/5", header_text, "Extracting DataPings", False))
        sleep(1)

        # 3. Extract DataPings
        data_pings = DataPings(filename, datagram, features)
        set_progress((60, "3/5", header_text, "Extracting DataSessions", False))

        # 4. Extract DataSessions
        data_session = DataSessions(pd.DataFrame([]), data_pings, features, 300, "")
        set_progress((80, "4/5", header_text, "Extracting Session Blocks", False))

        # 5. Extract Session Blocks
        data_session.extract_session_blocks()

        df_session = data_session.data.copy()
        df_pings = data_pings.data.copy()

        ident = driver.get_last_input("identifier")

        table = driver.get_df_from_db("identifier").copy()
        table.loc[len(table.index) - 1, "Type"] = "Feature"
        driver.df_to_sql_replace(table, "identifier")

        df_session["identifier"] = ident
        df_pings["identifier"] = ident

        driver.df_to_sql_append(df_session, "session")
        driver.filter_duplicates("session")

        driver.df_to_sql_append(df_pings, "pings")
        driver.filter_duplicates("pings")

        cluster_ids = data_session.get_cluster_ids()
        cluster_ids = cluster_ids.to_frame()
        cluster_ids["identifier"] = ident
        driver.df_to_sql_append(cluster_ids, "cluster_ids")

        driver.drop_current_table()
        set_progress((100, "5/5", header_text, "Loaded Data Successfully", False))

        return None, False, filename, dash.no_update
    return None, False, "", dash.no_update


@app.callback(
    Output(component_id="exportFunc", component_property="data"),
    Input(component_id="export", component_property="n_clicks"),
    State("overview_filename", component_property="children"),
    State("overview_lines", component_property="children"),
    State("overview_cal_days", component_property="children"),
    State("overview_metered_days", component_property="children"),
    State("select-date", "start_date"),
    State("select-date", "end_date"),
    State("graphs-store", "children"),
    State("additionals-store", "data"),
    State("license-store", "data"),
    prevent_initial_call=True,
)
def export_data(
    clicks: int,
    name: str,
    lines: str,
    cal_days: str,
    metered_days: str,
    start_date: str,
    end_date: str,
    graphs: list,
    additionals: dict,
    license: dict,
):
    """Export presentation on button click."""

    if clicks is not None:  # TODO: Do nothing if no data
        prs = Presentation("./assets/report_analysis_template.pptx")

        # title slide
        prs.slides[0].shapes[0].text = "Report Analysis: " + name
        prs.slides[0].shapes[1].text = start_date + " - " + end_date

        # data slide
        prs.slides[1].shapes[0].table.cell(1, 0).text = name
        prs.slides[1].shapes[0].table.cell(1, 1).text = cal_days
        prs.slides[1].shapes[0].table.cell(1, 2).text = metered_days
        prs.slides[1].shapes[0].table.cell(1, 3).text = lines

        # graph & statistic slides
        for option in DROPDOWN_OPTIONS:
            id = option["value"]
            name = option["label"]

            # TODO check if graph empty => layout 4
            slide = prs.slides.add_slide(
                prs.slide_layouts[2] if additionals[str(id)] else prs.slide_layouts[3]
            )
            slide.shapes.title.text = name

            # save and set graph
            graph_path = "./export/graphs/" + str(id) + ".png"
            Figure(graphs[id]["props"]["figure"]).write_image(graph_path)
            prsLib.set_graph(slide, graph_path)

            prsLib.set_table(slide, additionals[str(id)])

        # license usage slide
        if license:
            slide = prs.slides.add_slide(prs.slide_layouts[4])
            slide.shapes.title.text = "License Usage"
            print(license)
            prsLib.set_table(slide, license)
        
        prs.save("./export/report.pptx")


@app.long_callback(
    output=[Output(component_id="reset", component_property="n_clicks")],
    inputs=[Input(component_id="reset", component_property="n_clicks")],
    running=[
        (
            Output("reset_msg", "style"),
            {"top": "3%"},
            {"top": "-10%"},
        )
    ],
    # Uncomment the following line if the Database should NOT be resetted
    # after every reboot of the software
    # prevent_initial_call=True,
)
def reset_db(clicks: int):
    """
    Deletes the database
    Parameters
    ----------
    clicks : int
    Returns
    -------
    None
    """
    driver.drop_all()
    sleep(1.5)
    return dash.no_update


@app.callback(
    Output("ident", "value"), Input("ident", "value"), prevent_inital_call=True
)
def data_name_input(name):
    """
    Adds the input to the identifier table

    Parameters
    ----------
    name: String
        the name of the file identifier

    Returns
    -------
    None
        Resets the input field
    """
    if name is not None:
        driver.df_to_sql_append(
            pd.DataFrame({"FileIdentifier": [name], "Type": "unknown"}), "identifier"
        )
        driver.filter_duplicates("identifier")
    return None


@app.callback(
    Output("file-select", "options"),
    Output("file-select", "value"),
    Output("file-select-license", "options"),
    Output("file-select-license", "value"),
    Output("cluster_id-select", "options"),
    Output("cluster_id-select", "value"),
    Input("filename", "data"),
    Input("filename_license", "data"),
    Input("file-select", "value"),
    prevent_inital_call=True,
)
def set_select_options(filename: str, filename_license: str, file_select_value: str):
    """
    Update both select menus when uploading new files and update cluster_id select menu when another file is selected

    Parameter
    ---------
    filename : str
    file_select_value : str

    Returns
    -------
    list of str
        select option for file select
    str
        current selected file
    list of str
        select option for cluster_id select
    str
        current selected cluster_id
    """
    idents = []
    c_ids = []
    cur_c_id = dash.no_update
    if driver.check_if_table_exists("identifier"):
        idents = driver.get_df_from_db("identifier")
        if ctx.triggered_id == "filename_license":
            idents = (
                idents[idents["Type"] == "License"]["FileIdentifier"]
                .to_numpy(dtype=str)
                .flatten()
            )
            return (
                dash.no_update,
                dash.no_update,
                idents,
                idents[len(idents) - 1],
                dash.no_update,
                dash.no_update,
            )
        cur_c_id = "All Cluster-IDs"
        c_ids = driver.get_df_from_db("cluster_ids")
        if ctx.triggered_id == "filename":
            idents = (
                idents[idents["Type"] == "Feature"]["FileIdentifier"]
                .to_numpy(dtype=str)
                .flatten()
            )
            c_ids = c_ids[c_ids["identifier"] == idents[-1]]
        else:
            c_ids = c_ids[c_ids["identifier"] == file_select_value]
        c_ids = c_ids["cluster_id"].to_numpy().tolist()
        c_ids = ["All Cluster-IDs"] + c_ids
        if ctx.triggered_id == "filename":
            return idents, idents[-1], dash.no_update, dash.no_update, c_ids, cur_c_id
    return idents, dash.no_update, dash.no_update, dash.no_update, c_ids, cur_c_id


@app.callback(
    Output("settings-div", "style"),
    Input("open-settings-button", "n_clicks"),
    Input("close-settings-button", "n_clicks"),
)
def settings(open_settings, close_settings):
    triggered_id = ctx.triggered_id
    if triggered_id == "open-settings-button":
        return {"left": "0%"}
    else:
        return {"left": "-20%"}