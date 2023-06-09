import os
import shutil
from time import sleep
from typing import Callable

import dash_bootstrap_components as dbc
import dash_uploader as du
import diskcache
import pandas as pd
from dash import Dash, Input, Output, State, ctx, dash, dcc
from dash.long_callback import DiskcacheLongCallbackManager
from plotly.graph_objects import Figure
from pptx import Presentation
from pptx.util import Cm

import database.driver as driver
import vis.prs_lib as prs_lib
from computation.data import DataPings, DataSessions, LicenseUsage
from computation.features import Features
from dash_app import background, upload
from dash_app.upload import convert_report_to_df
from vis.additional_data_vis import get_license_usage_table
from vis.graph_vis import empty_fig
from vis.web_designs import DROPDOWN_OPTIONS, tab_layout

UPLOAD_CACHE_PATH = os.path.abspath("./cache/upload_data/")

# Diskcache - needed for long_callbacks
cache = diskcache.Cache(os.path.abspath("./cache"))
long_callback_manager = DiskcacheLongCallbackManager(cache)

# Dash
app = Dash(
    external_stylesheets=[dbc.themes.BOOTSTRAP],
    long_callback_manager=long_callback_manager,
    prevent_initial_callbacks=True,
)
app.title = "Token Dashboard Threedy"
app._favicon = "threedy_favicon.png"
app.layout = dbc.Container(tab_layout(), fluid=True)

du.configure_upload(app, UPLOAD_CACHE_PATH, use_upload_id=False)


@app.long_callback(
    output=[
        Output("dash-uploader", "isCompleted"),
        Output("filename", "data"),
        Output("filename_license", "data"),
    ],
    inputs=[
        Input("dash-uploader", "isCompleted"),
        State("dash-uploader", "fileNames"),
        Input("ident_num", "data"),
        State("ident_names", "data"),
    ],
    running=[
        (Output("main", "style"), {"filter": "grayscale(100%)"}, {}),
        (Output("header", "style"), {"filter": "grayscale(100%)"}, {}),
        (
            Output("progress_div", "style"),
            {"visibility": "visible"},
            {"visibility": "hidden"},
        ),
    ],
    progress=[
        Output("progress_bar", "value"),
        Output("progress_bar", "label"),
        Output("progress_bar_header", "children"),
        Output("progress_message", "children"),
        Output("modal_ident", "is_open"),
        Output("modal_header", "children"),
    ],
    prevent_inital_call=True,
)
def load_data(
    set_progress: Callable, is_com: bool, files: str, ident_num: int, ident_names
):
    """
    loads the data into the database

    Parameters
    ----------
    set_progress : Function to update the website while this callback function is running
    is_com : bool
        which indicates if the upload is completed
    files : String
        array of filenames (only first element used)
    ident_num : int
        number of already added identifier

    Returns
    -------
    bool which indicates if a download is complete
    str of the new feature file identifier
    str of the new license file identifier
    """

    """Upload file"""
    if is_com:
        header_text = "Upload Report"

        datagrams = convert_report_to_df(files[0])
        filename = files[0].split(".")[0]
        names = []
        for _, name in datagrams:
            names.append(name)

        if ident_num == -1:
            set_progress(
                (
                    100,
                    "0/5",
                    header_text,
                    "Waiting for input",
                    True,
                    "Please enter an identifier for the data from this file: "
                    + names[0],
                )
            )
        else:
            while ident_num < len(names):
                set_progress(
                    (
                        100,
                        "0/5",
                        header_text,
                        "Waiting for input",
                        True,
                        "Please enter an identifier for the data from this file: "
                        + names[ident_num],
                    )
                )

        return upload.prepare_data(
            set_progress, datagrams, filename, ident_num, ident_names
        )
    return dash.no_update, dash.no_update, dash.no_update


@app.callback(
    Output(component_id="report-statistics-table", component_property="children"),
    Output(component_id="graphs-store", component_property="children"),
    Output("additions-store", "data"),
    Output("select-date", "start_date"),
    Output("select-date", "end_date"),
    Input("select-date", "start_date"),
    Input("select-date", "end_date"),
    Input("filename", "data"),
    Input("graph-type", "value"),
    Input("file-select-feature", "value"),
    Input("cluster_id-select", "value"),
    Input("multi_cluster", "value"),
    Input("time-reset", "n_clicks"),
    Input(component_id="apply-report-selection", component_property="n_clicks"),
    prevent_inital_call=True,
)
def update_output_div(
    start_date: str,
    end_date: str,
    filename: str,
    graph_type: str,
    file_select_value: list,
    c_id_select: str,
    multi_cluster: str,
    time_reset: int,
    clicks: int,
):
    """
    Parameters
    ----------
    start_date : String which represents the selected entry of the beginning day in the calendar tool
                    with the id 'select-date'
    end_date : String which represents the selected entry of the ending day in the calendar tool
                    with the id 'select-date'
    filename : str
    graph_type : String, either "bar", "line" or "auto"
    file_select_value: list of String
        the selected file identifier
    c_id_select: String
        the selected cluster id
    multi_cluster: bool
        True if cluster_id should be aggregated over all file identifier
    time_reset: int
        Only used for updates, indicates if the time interval should be maximised
    clicks : int
        clicks of apply report selection button

    main computation of the frontend

    Returns
    -------
    String which the number of days inbetween the selected days of the calendar tool
            or the first and last day of the dataframe
    String which represents the metered days in the dataframe
    dcc.Graph which represents the graph with the id 'graph1'
    dcc.Graph which represents the graph with the id 'graph2'
    String which represents additional data to the graph with the id 'graph1'
    String which represents additional data to the graph with the id 'graph2'
    date.Date which represents the first selected day in the calendar tool
    date.Date which represents the last selected day in the calendar tool
    """
    # if no file is selected or if selection hasn't been applied yet
    if ctx.triggered_id == "file-select-feature" or not file_select_value:
        return (
            dash.no_update,  # overview_table data
            dash.no_update,  # overview_table data
            dash.no_update,  # overview_table data
            dash.no_update,  # overview_table data
            dash.no_update,  # export data
            dash.no_update,  # export data
            dash.no_update,  # time interval data
            dash.no_update,  # time interval data
        )
    if driver.check_if_table_exists("session"):
        features = Features().get_data_features()

        """set current cluster id"""
        c_id = None
        # if not (c_id_select == "All Cluster-IDs" or c_id_select == []):
        if not c_id_select == "All Cluster-IDs":
            c_id = c_id_select

        """converting the dict containing all data back into a pd.Dataframe"""
        sql_session = driver.get_df_from_db("session")
        """create DataSessions object"""
        data_pings = DataPings(filename, driver.get_df_from_db("pings"), features, c_id)
        sessions = DataSessions(
            sql_session, data_pings, features, 300, file_select_value, c_id
        )

        """checking if new data is loaded and new initial dates should be set"""
        new_data = False
        if ctx.triggered_id == "filename" or ctx.triggered_id == "time-reset":
            new_data = True

        """setting the first and last dates for the represented data"""
        first_date = background.select_date(start_date, sessions.data, True, new_data)
        last_date = background.select_date(end_date, sessions.data, False, new_data)

        """trim dataframe to selected time interval"""
        sessions.crop_data(first_date, last_date)

        empty_val = False
        if (
            len(
                sessions.data[sessions.data["identifier"].isin(file_select_value)].index
            )
            == 0
        ):
            empty_val = True

        """store data of current settings in database"""
        if (graph_type == "automatic") and (not empty_val):
            graph_type = "bar" if (len(data_pings.get_metered_days()) <= 2) else "line"

        graphs = []
        additions = dict()
        if "Aggregate Cluster ID data" in multi_cluster:
            multi_cluster_bool = True
        else:
            multi_cluster_bool = False
        for option in DROPDOWN_OPTIONS:
            dropdown_id = option["label"]
            if empty_val:
                fig = empty_fig()
                additional = pd.DataFrame()
            else:
                fig, additional = background.select_graph(
                    dropdown_id,
                    sessions,
                    file_select_value,
                    graph_type,
                    multi_cluster_bool,
                )
            graphs.append(dcc.Graph(figure=fig, className="graph"))
            additions[str(option["value"])] = additional.to_dict()

        """get values for file statistics"""
        report_statistics_table = background.get_report_statistics_table()
    else:
        """get values for file statistics"""
        report_statistics_table = background.get_report_statistics_table()

        """no data -> no updates for visuals"""
        graphs = dash.no_update
        additions = dash.no_update

        """set first and last date"""
        first_date = start_date
        last_date = end_date

    return (
        report_statistics_table,  # report-statistics-table
        graphs,  # export data
        additions,  # export data
        first_date,  # time interval data
        last_date,  # time interval data
    )


@app.callback(
    Output(component_id="graph_data3", component_property="children"),
    Output("license-store", "data"),
    Input("filename_license", "data"),
    prevent_inital_call=True,
)
def update_output_license(filename: str):
    """
    Parameters
    ----------
    filename: String

    main computation of the frontend

    Returns
    -------
    dcc.Graph which represents the graph with the id 'graph3'
    """
    if driver.check_if_table_exists("license"):
        license_data = background.get_license_data()
        license_usage = LicenseUsage(license_data)
        additional = get_license_usage_table(
            license_usage, background.get_license_identifier()
        )
        return (
            dbc.Table.from_dataframe(additional, style={"text-align": "right"}),
            additional.to_dict(),
        )
    else:
        return dash.no_update, dash.no_update


@app.callback(
    Output(component_id="graph1", component_property="children"),
    Output(component_id="graph2", component_property="children"),
    Output(component_id="graph_data1", component_property="children"),
    Output(component_id="graph_data2", component_property="children"),
    Input("graphs-store", "children"),
    Input("additions-store", "data"),
    Input(component_id="dropdown1", component_property="value"),
    Input(component_id="dropdown2", component_property="value"),
    prevent_inital_call=True,
)
def update_dropdown(
    graphs: list,
    additions: dict,
    drop1: int,
    drop2: int,
):
    if not graphs:
        graph1 = empty_fig()
        graph2 = empty_fig()
        additional1 = ""
        additional2 = ""
    else:
        graph1 = graphs[drop1]
        additional1 = dbc.Table.from_dataframe(
            pd.DataFrame.from_dict(additions[str(drop1)]), style={"text-align": "right"}
        )

        graph2 = graphs[drop2]
        additional2 = dbc.Table.from_dataframe(
            pd.DataFrame.from_dict(additions[str(drop2)]), style={"text-align": "right"}
        )
    return graph1, graph2, additional1, additional2


@app.callback(
    Output(component_id="exportFunc", component_property="data"),
    Input(component_id="export", component_property="n_clicks"),
    State("select-date", "start_date"),
    State("select-date", "end_date"),
    State("graphs-store", "children"),
    State("additions-store", "data"),
    State("license-store", "data"),
    prevent_initial_call=True,
)
def export_data(
    clicks: int,
    start_date: str,
    end_date: str,
    graphs: list,
    additions: dict,
    license_data: dict,
):
    """Export presentation on button click."""

    if clicks is not None:
        prs = Presentation("./assets/report_analysis_template.pptx")

        # title slide
        prs.slides[0].shapes[0].text = "Report Analysis"
        prs.slides[0].shapes[1].text = start_date + " - " + end_date

        # report Statistics slide
        if driver.check_if_table_exists("report_statistics"):
            report_statistics = driver.get_df_from_db("report_statistics")
            slide = prs.slides.add_slide(prs.slide_layouts[4])
            slide.shapes.title.text = "Report Statistics"
            prs_lib.set_table(slide, report_statistics.to_dict(), Cm(5))

        # graph & statistic slides
        for option in DROPDOWN_OPTIONS:
            dropdown_id = option["value"]
            name = option["label"]

            slide = prs.slides.add_slide(
                prs.slide_layouts[2]
                if additions[str(dropdown_id)]
                else prs.slide_layouts[3]
            )
            slide.shapes.title.text = name

            # save and set graph
            graph_path = "./export/graphs/" + str(dropdown_id) + ".png"
            Figure(graphs[dropdown_id]["props"]["figure"]).write_image(graph_path)
            prs_lib.set_graph(slide, graph_path)

            prs_lib.set_table(slide, additions[str(dropdown_id)])

        # license usage slide
        if license_data:
            slide = prs.slides.add_slide(prs.slide_layouts[4])
            slide.shapes.title.text = "License Usage"
            prs_lib.set_table(slide, license_data)

        prs.save("./export/report.pptx")
        return dcc.send_file("./export/report.pptx")


@app.long_callback(
    output=[Output(component_id="reset", component_property="n_clicks")],
    inputs=[Input(component_id="reset", component_property="n_clicks")],
    running=[
        (
            Output("reset_msg", "style"),
            {"top": "3%"},
            {"top": "-10%"},
        )
    ],
    # Uncomment the following line if the Database should NOT be resetted
    # after every reboot of the software
    prevent_initial_call=True,
)
def reset_db(clicks: int):
    """
    Deletes the database

    Parameters
    ----------
    clicks : int

    Returns
    -------
    dash.no_update
    """
    driver.drop_all()
    shutil.rmtree(UPLOAD_CACHE_PATH)
    sleep(1.5)
    return dash.no_update


@app.callback(
    Output("ident", "value"),
    Output("ident_num", "data"),
    Output("confirm", "n_clicks"),
    Output("ident_names", "data"),
    Input("confirm", "n_clicks"),
    Input("dash-uploader", "fileNames"),
    State("ident", "value"),
    State("ident_num", "data"),
    State("all_file_check", "value"),
    State("ident_names", "data"),
    prevent_inital_call=True,
)
def data_name_input(confirm, file, name, num, checkbox, ident_names):
    """
    Adds the input to the identifier table

    Parameters
    ----------
    confirm: int
        shows if the confirm button was clicked
    file: str
        not used, only for updates
    name: String
        the name of the file identifier
    num: int
        The current amount of file identifier used for a compressed file
    checkbox: list(str)
        information if the identifier should be used for all files in a compressed

    Returns
    -------
    None
        Resets the input field
    int
        Number of already set identifier
    None
        Resets the confirm button
    """
    if ctx.triggered_id == "dash-uploader":
        return dash.no_update, 0, dash.no_update, None

    if ident_names is None:
        name_lst = []
    else:
        name_lst = ident_names

    if name is not None:
        while confirm is None:
            sleep(0.1)
        driver.df_to_sql_append(
            pd.DataFrame({"FileIdentifier": [name], "Type": "unknown"}), "identifier"
        )
        driver.filter_duplicates("identifier", identifier=["FileIdentifier"])
        name_lst.append(name)
        if "Use Identifier for all files" in checkbox:
            return None, -1, None, name_lst
        return None, num + 1, None, name_lst
    return dash.no_update, dash.no_update, dash.no_update, dash.no_update


@app.callback(
    Output("file-select-feature", "options"),
    Output("file-select-feature", "value"),
    Output("cluster_id-select", "options"),
    Output("cluster_id-select", "value"),
    Input("filename", "data"),
    Input("file-select-feature", "value"),
    Input(component_id="apply-report-selection", component_property="n_clicks"),
    prevent_inital_call=True,
)
def set_select_options(filename: str, file_select_value: str, clicks: int):
    """
    Update select menus of cluster_ids, license_id and feature_id when a new file gets uploaded
    or another cluster_id or
     feature_id is selected

    Parameter
    ---------
    filename : str
    file_select_value : list of str
    clicks : int
        clicks of apply report selection button

    Returns
    -------
    list of str
        select option for feature file select
    list of str
        current selected feature files
    list of str
        select option for license file select
    str
        current selected license file
    list of str
        select option for cluster_id select
    str
        current selected cluster_id
    """
    sleep(0.5)
    if (
        driver.check_if_table_exists("identifier")
        and not ctx.triggered_id == "file-select-feature"
        and not (ctx.triggered_id == "apply-report-selection" and not file_select_value)
    ):
        idents = background.get_feature_identifier()

        cur_c_id = "All Cluster-IDs"
        c_ids = driver.get_df_from_db("cluster_ids")
        if ctx.triggered_id == "filename":
            c_ids = c_ids[c_ids["identifier"] == idents[-1]]
        else:
            c_ids = c_ids[c_ids["identifier"].isin(file_select_value)]
        c_ids = c_ids["cluster_id"].drop_duplicates()
        c_ids = c_ids.to_numpy().tolist()
        c_ids = ["All Cluster-IDs"] + c_ids
        if ctx.triggered_id == "filename":
            return idents, [idents[-1]], c_ids, cur_c_id
        return dash.no_update, dash.no_update, c_ids, cur_c_id
    return dash.no_update, dash.no_update, dash.no_update, dash.no_update


@app.callback(
    Output("settings-div", "style"),
    Input("open-settings-button", "n_clicks"),
    Input("close-settings-button", "n_clicks"),
    prevent_inital_call=True,
)
def settings(open_settings, close_settings):
    triggered_id = ctx.triggered_id

    if triggered_id == "open-settings-button":
        if close_settings is None:
            if open_settings % 2:
                return {"left": "0%"}
            else:
                return {"left": "-20%"}
        elif (open_settings + close_settings) % 2:
            return {"left": "0%"}
        else:
            return {"left": "-20%"}

    if triggered_id == "close-settings-button":
        return {"left": "-20%"}

    return {"left": "-20%"}
